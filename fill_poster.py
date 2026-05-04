#!/usr/bin/env python3
"""Fill SmartCareer AI graduation project poster."""
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
import copy

STUDENT_ID = "210004"
SUPERVISOR = "Farhod Mahmudxo'jayev"
PROFESSOR = "Eugene Castro"
PROFESSOR_EMAIL = "e.castro@centralasian.uz"

def set_text(shape, text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        tf.paragraphs[i]._p.getparent().remove(tf.paragraphs[i]._p)
    para = tf.paragraphs[0]
    for run in list(para.runs):
        run._r.getparent().remove(run._r)
    lines = text.split('\n')
    for li, line in enumerate(lines):
        if li == 0:
            r = para.add_run()
            r.text = line
        else:
            new_p = copy.deepcopy(para._p)
            for old_r in new_p.findall(qn('a:r')):
                new_p.remove(old_r)
            new_r = copy.deepcopy(para._p.findall(qn('a:r'))[0])
            new_r.find(qn('a:t')).text = line
            new_p.append(new_r)
            tf._txBody.append(new_p)

def find(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

prs = Presentation('StudentIDs_Poster.pptx')
slide = prs.slides[0]

# INTRODUCTION
s = find(slide, "TextBox 31")
if s:
    set_text(s, (
        "INTRODUCTION\n\n"
        "The transition from education to employment is one of the most critical challenges "
        "faced by university students and recent graduates. Despite the abundance of online job "
        "portals, many young job seekers struggle to create professional, ATS-optimized resumes "
        "and effectively match their skills to available opportunities. Simultaneously, employers "
        "face difficulties in efficiently filtering and identifying the most suitable candidates "
        "from large applicant pools. This gap highlights the need for an intelligent, AI-driven "
        "platform that can streamline the entire career development pipeline — from resume creation "
        "to job matching and application tracking. SmartCareer AI addresses this challenge by "
        "providing an integrated platform that leverages artificial intelligence (Google Gemini "
        "and OpenAI GPT-4) to automate resume generation, analyze job-resume compatibility, "
        "generate tailored cover letters, and facilitate intelligent job matching, thereby "
        "empowering students to compete more effectively in today's job market."
    ))
    print("  OK Introduction")

# OBJECTIVES
s = find(slide, "TextBox 37")
if s:
    set_text(s, (
        "PROJECT OBJECTIVES\n\n"
        "1. To develop an AI-powered resume builder that generates professional, ATS-optimized "
        "resumes using OpenAI GPT-4 and Google Gemini, reducing manual effort by ~70%.\n"
        "2. To implement an intelligent job matching algorithm that analyzes candidate profiles "
        "against job requirements and provides compatibility scores with actionable recommendations.\n"
        "3. To build a complete multi-role web platform (Student, Company/HR, Admin) with secure "
        "authentication, application tracking, and real-time notifications.\n"
        "4. To integrate AI-driven cover letter generation that tailors content based on the "
        "candidate's resume and specific job descriptions.\n"
        "5. To design a scalable, production-ready architecture using modern frameworks (FastAPI "
        "+ Next.js) with CI/CD pipelines, monitoring, and deployment automation."
    ))
    print("  OK Objectives")

# METHODOLOGY
s = find(slide, "TextBox 35")
if s:
    set_text(s, (
        "METHODOLOGY\n\n"
        "The development followed an Agile methodology with iterative sprints:\n\n"
        "Backend: Built with FastAPI (Python), implementing RESTful APIs with versioned endpoints. "
        "SQLAlchemy ORM handles database operations with PostgreSQL. Alembic manages migrations. "
        "JWT-based authentication with OAuth2 (Google) ensures secure access.\n\n"
        "AI Layer: A dual-provider service supports Google Gemini (primary, free) and OpenAI GPT-4 "
        "(fallback) with automatic failover. Structured prompts and JSON-validated responses ensure "
        "consistent AI output quality.\n\n"
        "Frontend: Next.js 14 with React 18 and TypeScript provides server-side rendering. "
        "Zustand manages state, Tailwind CSS and Radix UI deliver a modern, accessible interface.\n\n"
        "Testing & CI/CD: Pytest for backend, Playwright for E2E, GitHub Actions for CI/CD."
    ))
    print("  OK Methodology")

# RESULTS
s = find(slide, "TextBox 36")
if s:
    set_text(s, (
        "RESULTS\n\n"
        "The AI Resume Builder successfully generates tailored, ATS-optimized resumes using dual "
        "AI providers (Gemini/GPT-4), reducing manual resume creation effort by approximately 70%. "
        "The intelligent Job Matching system analyzes candidate profiles against job requirements, "
        "providing match scores (0-100), identifying skill gaps, and offering actionable recommendations. "
        "The platform delivers a complete 3-role ecosystem: Students build resumes, search/apply for "
        "jobs, and track applications; Companies post jobs, manage applicants, and filter candidates "
        "with AI assistance; Admins have a comprehensive dashboard for system monitoring. The Cover "
        "Letter Generator produces personalized, job-specific cover letters. The system architecture "
        "proved robust under testing, with Sentry integration for error monitoring, Redis-based rate "
        "limiting, and GZip compression ensuring production-grade performance."
    ))
    print("  OK Results")

# CONCLUSION
s = find(slide, "TextBox 50")
if s:
    set_text(s, (
        "CONCLUSION\n\n"
        "This project has successfully achieved all five objectives, delivering a fully functional "
        "AI-powered career platform that bridges the gap between job seekers and employers. "
        "SmartCareer AI has evolved beyond a simple job portal into an intelligent career development "
        "tool that leverages cutting-edge AI capabilities. The dual AI provider architecture "
        "(Gemini + OpenAI) ensures reliability and cost-effectiveness, while the modular codebase "
        "and comprehensive testing infrastructure provide a solid foundation for future development. "
        "Key lessons: modular architecture for managing full-stack solo development, iterative prompt "
        "engineering for consistent AI quality, and automated testing/CI/CD for code reliability."
    ))
    print("  OK Conclusion")

# FUTURE WORK
s = find(slide, "TextBox 52")
if s:
    set_text(s, (
        "RECOMMENDATION AND FUTURE WORK\n\n"
        "1. Mobile Application: Develop a React Native mobile app for on-the-go job searching.\n"
        "2. AI Interview Coach: Implement AI-powered mock interviews with real-time feedback.\n"
        "3. ML Job Recommendations: Train models on user behavior for personalized recommendations.\n"
        "4. University Partnership: Integrate with career centers for grant/scholarship search.\n"
        "5. Localization: Expand Uzbek and Russian language support for Central Asian markets.\n"
        "6. Payment Integration: Integrate Click.uz and Payme alongside existing Stripe."
    ))
    print("  OK Future Work")

# SUPERVISOR
s = find(slide, "TextBox 53")
if s:
    set_text(s, (
        f"SUPERVISOR\n\n"
        f"Mr. {SUPERVISOR}\n\n"
        f"PROFESSOR\n\n"
        f"{PROFESSOR}\n"
        f"{PROFESSOR_EMAIL}"
    ))
    print("  OK Supervisor")

# PROJECT OUTPUT
s = find(slide, "TextBox 57")
if s:
    set_text(s, (
        "PROJECT OUTPUT\n\n"
        "The following screenshots show the SmartCareer AI web platform interface, including "
        "the AI Resume Builder, Job Matching dashboard, Application Tracker, Company HR portal, "
        "and Admin dashboard — demonstrating all three user roles (Student, Company, Admin)."
    ))
    print("  OK Project Output")

# ACKNOWLEDGMENT
s = find(slide, "TextBox 49")
if s:
    set_text(s, (
        "ACKNOWLEDGMENT\n\n"
        f"I would like to express my sincere gratitude to my supervisor, Mr. {SUPERVISOR}, "
        f"and to Professor {PROFESSOR} for their invaluable guidance, constructive feedback, "
        "and continuous support throughout this project. My sincere appreciation also goes to "
        "the Engineering School and the Central Asian University community for their continuous "
        "academic inspiration and the enriching environment they provide for learning and innovation."
    ))
    print("  OK Acknowledgment")

out = f'{STUDENT_ID}_Poster.pptx'
prs.save(out)
print(f"\nSaved: {out}")
