#!/usr/bin/env python3
"""Fix issues found in both generated presentations."""
from pptx import Presentation
from pptx.oxml.ns import qn
import copy

STUDENT_ID = "210004"
STUDENT_NAME = "Djumabaev Abdurashid"
PROJECT_TITLE = "SmartCareer AI"
PROJECT_SUBTITLE = "An AI-Powered Career Platform for Students and Employers"

def set_text(shape, text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        tf.paragraphs[i]._p.getparent().remove(tf.paragraphs[i]._p)
    para = tf.paragraphs[0]
    for run in list(para.runs):
        run._r.getparent().remove(run._r)
    lines = text.split('\n')
    for li, line in enumerate(lines):
        if li == 0:
            r = para.add_run()
            r.text = line
        else:
            new_p = copy.deepcopy(para._p)
            for old_r in new_p.findall(qn('a:r')):
                new_p.remove(old_r)
            new_r = copy.deepcopy(para._p.findall(qn('a:r'))[0])
            new_r.find(qn('a:t')).text = line
            new_p.append(new_r)
            tf._txBody.append(new_p)

# =========================================================================
# FIX POSTER
# =========================================================================
print("FIXING POSTER...")
prs = Presentation(f'{STUDENT_ID}_Poster.pptx')
slide = prs.slides[0]

for shape in slide.shapes:
    if shape.shape_type == 6:  # GROUP
        for child in shape.shapes:
            if child.name == "TextBox 25":
                set_text(child, PROJECT_TITLE)
                print(f"  FIXED: Project Title in poster header -> '{PROJECT_TITLE}'")
            elif child.name == "TextBox 26":
                set_text(child, PROJECT_SUBTITLE)
                print(f"  FIXED: Project Subtitle -> '{PROJECT_SUBTITLE}'")
            elif child.name == "TextBox 29":
                set_text(child, STUDENT_NAME)
                print(f"  FIXED: Student Name -> '{STUDENT_NAME}'")
            elif child.name == "TextBox 30":
                set_text(child, f"{STUDENT_ID} - BSc in Computer Science")
                print(f"  FIXED: Student ID line -> '{STUDENT_ID} - BSc in Computer Science'")

prs.save(f'{STUDENT_ID}_Poster.pptx')
print(f"  Saved: {STUDENT_ID}_Poster.pptx")

# =========================================================================
# FIX PRESENTATION
# =========================================================================
print("\nFIXING PRESENTATION...")
prs = Presentation(f'{STUDENT_ID}_SmartCareer_AI.pptx')

# Fix Slide 24: "Furture Scope" -> "Future Scope"
slide = prs.slides[23]
for shape in slide.shapes:
    if shape.name == "TextBox 3" and hasattr(shape, 'text'):
        if "Furture" in shape.text:
            set_text(shape, "Future Scope")
            print("  FIXED: Slide 24 typo 'Furture Scope' -> 'Future Scope'")

# Fix Slide 27: Remove empty rows from contribution table
slide = prs.slides[26]
for shape in slide.shapes:
    if shape.shape_type == 19:  # TABLE
        table = shape.table
        # Mark empty rows with a dash to look cleaner
        for ri in range(2, len(table.rows)):
            row = table.rows[ri]
            all_empty = all(cell.text.strip() == "" for cell in row.cells)
            if all_empty:
                row.cells[0].text = "—"
                row.cells[1].text = "—"
                row.cells[2].text = "—"
                row.cells[3].text = "—"
                row.cells[4].text = "—"
        print("  FIXED: Slide 27 empty table rows marked with dashes")

prs.save(f'{STUDENT_ID}_SmartCareer_AI.pptx')
print(f"  Saved: {STUDENT_ID}_SmartCareer_AI.pptx")

# =========================================================================
# FINAL VERIFICATION
# =========================================================================
print("\n" + "=" * 60)
print("FINAL VERIFICATION")
print("=" * 60)

# Verify poster
prs = Presentation(f'{STUDENT_ID}_Poster.pptx')
slide = prs.slides[0]
print("\nPOSTER HEADER CHECK:")
for shape in slide.shapes:
    if shape.shape_type == 6:
        for child in shape.shapes:
            if hasattr(child, 'text') and child.text.strip():
                print(f"  {child.name}: {repr(child.text[:80])}")

# Verify presentation
prs = Presentation(f'{STUDENT_ID}_SmartCareer_AI.pptx')
print("\nPRESENTATION SLIDE-BY-SLIDE CHECK:")
issues = []
for i, slide in enumerate(prs.slides):
    has_content = False
    slide_texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            t = shape.text.strip()
            if t not in [str(i+1)]:  # Skip slide numbers
                has_content = True
                slide_texts.append(t[:40])
        if shape.shape_type == 19:
            has_content = True
    
    # Check for template placeholder text
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            txt = shape.text.lower()
            if any(p in txt for p in ['item\nitem', '[project title]', '[name of', '[date of', '[dr.', 'student ids', 'you should', 'describe how', 'provide a visual', 'present here', 'show a clear', 'state here', 'show the most', '"in this slide']):
                issues.append(f"Slide {i+1}: Possible unfilled placeholder in {shape.name}: {repr(shape.text[:60])}")
    
    status = "OK" if has_content else "EMPTY?"
    title = slide_texts[0] if slide_texts else "(no text)"
    print(f"  Slide {i+1}: [{status}] {title}")

if issues:
    print(f"\n⚠️  POTENTIAL ISSUES FOUND ({len(issues)}):")
    for issue in issues:
        print(f"  - {issue}")
else:
    print("\n✅ No unfilled placeholders detected!")
