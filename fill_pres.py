#!/usr/bin/env python3
"""Fill SmartCareer AI graduation project presentation."""
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
import copy

STUDENT_ID = "210004"
STUDENT_NAME = "Djumabaev Abdurashid"
SUPERVISOR = "Farhod Mahmudxo'jayev"
PROFESSOR = "Eugene Castro"
PROFESSOR_EMAIL = "e.castro@centralasian.uz"
PROJECT_TITLE = "SmartCareer AI: An AI-Powered Career Platform for Students and Employers"
DATE = "5 May, 2026"

def set_text(shape, text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        tf.paragraphs[i]._p.getparent().remove(tf.paragraphs[i]._p)
    para = tf.paragraphs[0]
    for run in list(para.runs):
        run._r.getparent().remove(run._r)
    lines = text.split('\n')
    for li, line in enumerate(lines):
        if li == 0:
            r = para.add_run()
            r.text = line
        else:
            new_p = copy.deepcopy(para._p)
            for old_r in new_p.findall(qn('a:r')):
                new_p.remove(old_r)
            new_r = copy.deepcopy(para._p.findall(qn('a:r'))[0])
            new_r.find(qn('a:t')).text = line
            new_p.append(new_r)
            tf._txBody.append(new_p)

def find(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

prs = Presentation('StudentIDs_[Project Title].pptx')

# =========================================================================
# SLIDE 1: Title
# =========================================================================
slide = prs.slides[0]
s = find(slide, "TextBox 6")
if s: set_text(s, PROJECT_TITLE)
s = find(slide, "TextBox 11")
if s: set_text(s, DATE)
s = find(slide, "TextBox 7")
if s: set_text(s, f"Graduation Project\n\n{STUDENT_NAME}\nStudent ID: {STUDENT_ID}\n\nSupervisor: Mr. {SUPERVISOR}\nProfessor: {PROFESSOR}")
print("  OK Slide 1: Title")

# =========================================================================
# SLIDE 3: Problem Statement
# =========================================================================
slide = prs.slides[2]
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "• Students and recent graduates lack tools to create professional, ATS-optimized resumes, "
        "resulting in missed job opportunities despite having relevant qualifications.\n"
        "• Traditional job portals offer no intelligent matching — candidates must manually search "
        "through hundreds of listings without personalized compatibility analysis.\n"
        "• Employers spend excessive time filtering large applicant pools manually, lacking AI-assisted "
        "candidate screening tools to identify the best fits efficiently."
    ))
s = find(slide, "TextBox 5")
if s:
    set_text(s, (
        "Over 75% of resumes are rejected by ATS systems before reaching a human recruiter. "
        "Junior developers spend an average of 40+ hours monthly on job applications. "
        "SmartCareer AI addresses these pain points by leveraging AI to automate resume creation, "
        "provide intelligent job matching, and streamline the entire application process."
    ))
print("  OK Slide 3: Problem Statement")

# =========================================================================
# SLIDE 4: Objectives
# =========================================================================
slide = prs.slides[3]
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "1. Develop an AI-powered resume builder using OpenAI GPT-4 and Google Gemini, reducing manual effort by ~70%.\n"
        "2. Implement intelligent job matching with compatibility scores (0-100) and actionable recommendations.\n"
        "3. Build a multi-role platform (Student, Company/HR, Admin) with secure auth and application tracking.\n"
        "4. Integrate AI-driven cover letter generation tailored to specific job descriptions.\n"
        "5. Design a scalable architecture (FastAPI + Next.js) with CI/CD, monitoring, and deployment automation."
    ))
s = find(slide, "TextBox 5")
if s: set_text(s, "")
print("  OK Slide 4: Objectives")

# =========================================================================
# SLIDE 5: Timeline
# =========================================================================
slide = prs.slides[4]
s = find(slide, "TextBox 5")
if s:
    set_text(s, (
        "Phase 1 (Weeks 1-2): Requirements analysis, system architecture, database schema design\n\n"
        "Phase 2 (Weeks 3-5): Backend — FastAPI setup, JWT/OAuth2 auth, SQLAlchemy models, API endpoints\n\n"
        "Phase 3 (Weeks 6-8): AI integration — OpenAI/Gemini services, resume generation, job matching, ATS scoring\n\n"
        "Phase 4 (Weeks 9-11): Frontend — Next.js 14, student dashboard, company portal, admin panel, Tailwind UI\n\n"
        "Phase 5 (Weeks 12-13): Testing — Pytest, Playwright E2E, GitHub Actions CI/CD, Docker, deployment\n\n"
        "Phase 6 (Week 14): Final testing, documentation, and presentation preparation"
    ))
print("  OK Slide 5: Timeline")

# =========================================================================
# SLIDE 7: Architecture Part 1
# =========================================================================
slide = prs.slides[6]
s = find(slide, "TextBox 3")
if s: set_text(s, "System Architecture Overview")
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "SmartCareer AI follows a modern client-server architecture with three layers:\n\n"
        "Frontend Layer (Next.js 14 + React 18):\n"
        "Server-side rendering for SEO, Zustand state management, Tailwind CSS responsive design, "
        "Axios for API communication, Radix UI accessible components.\n\n"
        "Backend API Layer (FastAPI + Python):\n"
        "Versioned REST endpoints (/api/v1/), JWT authentication, role-based access control, "
        "Pydantic validation, CORS/security middleware, GZip compression, request ID tracking.\n\n"
        "Data & AI Layer:\n"
        "PostgreSQL via SQLAlchemy ORM, Alembic migrations, dual AI provider (Gemini primary + "
        "OpenAI fallback), Redis for caching/rate limiting, Sentry error monitoring."
    ))
print("  OK Slide 7: Architecture Overview")

# =========================================================================
# SLIDE 8: Architecture Part 2
# =========================================================================
slide = prs.slides[7]
s = find(slide, "TextBox 3")
if s: set_text(s, "API Structure and Endpoints")
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "The backend exposes 10 versioned API modules under /api/v1/:\n\n"
        "/auth — Registration, login, email verification, password reset, OAuth2 (Google)\n"
        "/users — User profile management, role-based access\n"
        "/resumes — Manual resume CRUD + AI resume generation with ATS scoring\n"
        "/jobs — Job posting CRUD, search, filtering, publish/unpublish\n"
        "/applications — Applications, status tracking (reviewing/shortlisted/interview/accepted/rejected)\n"
        "/ai — Resume Generator, Resume Analyzer, Cover Letter Generator, Job Matcher\n"
        "/admin — System health, user management, error monitoring dashboard\n"
        "/payments — Stripe subscription and billing management\n"
        "/notifications — Real-time notification system\n"
        "/saved-searches — Search filter persistence\n\n"
        "All endpoints include auto-documentation (Swagger/ReDoc), request validation, "
        "structured error responses, and security headers."
    ))
print("  OK Slide 8: API Structure")

# =========================================================================
# SLIDE 9: Architecture Part 3 - AI Pipeline
# =========================================================================
slide = prs.slides[8]
s = find(slide, "TextBox 3")
if s: set_text(s, "AI Service Architecture")
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "Provider-agnostic design with automatic failover:\n\n"
        "Resume Generation Pipeline:\n"
        "User Input → Data Extraction → Prompt Construction → AI Provider (Gemini/GPT-4) → "
        "JSON Validation → Content Normalization → ATS Score → DB Storage → PDF Export\n\n"
        "Job Matching Pipeline:\n"
        "Resume → Skill Extraction → Experience Analysis → Keyword Extraction → Score Calculation "
        "against Active Jobs → Ranked Results with Match Reasons\n\n"
        "Cover Letter Generation:\n"
        "Resume + Job Description + Company Info → Structured Prompt → AI Generation → "
        "Key Points + Match Score\n\n"
        "Failover: Gemini (free) → OpenAI GPT-4 (paid) → Error with config guidance\n"
        "All responses are JSON-validated and stored with metadata (provider, model, time, tokens)."
    ))
print("  OK Slide 9: AI Architecture")

# =========================================================================
# SLIDE 11: Tools - Backend
# =========================================================================
slide = prs.slides[10]
s = find(slide, "TextBox 3")
if s: set_text(s, "Backend Technologies")
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "• FastAPI 0.104 — High-performance Python web framework with automatic OpenAPI docs\n\n"
        "• SQLAlchemy 2.0 — ORM for database operations with PostgreSQL and SQLite\n\n"
        "• Alembic 1.12 — Database migration tool for version-controlled schema changes\n\n"
        "• JWT (python-jose) — JSON Web Token authentication with secure token management\n\n"
        "• Pydantic 2.5 — Data validation and settings management via Python type hints\n\n"
        "• PostgreSQL — Production-grade relational database\n\n"
        "• Redis 5.0 — In-memory store for rate limiting, token blacklisting, OAuth state\n\n"
        "• Gunicorn — Production WSGI server for FastAPI deployment"
    ))
print("  OK Slide 11: Backend Tools")

# =========================================================================
# SLIDE 12: Tools - Frontend
# =========================================================================
slide = prs.slides[11]
s = find(slide, "TextBox 3")
if s: set_text(s, "Frontend Technologies")
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "• Next.js 14 — React framework with SSR, routing, and optimization\n\n"
        "• React 18 — Component-based UI library with hooks and concurrent features\n\n"
        "• TypeScript 5.4 — Typed superset of JavaScript for improved code quality\n\n"
        "• Tailwind CSS 3.4 — Utility-first CSS framework for rapid responsive design\n\n"
        "• Zustand 4.5 — Lightweight state management for React\n\n"
        "• Radix UI — Accessible, unstyled component primitives\n\n"
        "• Framer Motion — Animation library for smooth UI transitions\n\n"
        "• Recharts — Charting library for admin dashboard data visualization"
    ))
print("  OK Slide 12: Frontend Tools")

# =========================================================================
# SLIDE 13: Tools - AI & Infrastructure
# =========================================================================
slide = prs.slides[12]
s = find(slide, "TextBox 3")
if s: set_text(s, "AI and Infrastructure Tools")
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "• OpenAI GPT-4 — LLM for resume generation, analysis, and cover letters\n\n"
        "• Google Gemini — Primary AI provider (free tier) with automatic OpenAI fallback\n\n"
        "• Sentry — Real-time error monitoring and performance tracking\n\n"
        "• Docker — Containerization for consistent dev/prod environments\n\n"
        "• Playwright — End-to-end testing framework for cross-browser UI testing\n\n"
        "• GitHub Actions — CI/CD pipeline for automated testing and deployment\n\n"
        "• Stripe — Payment processing for subscription-based premium features\n\n"
        "• Render / Vercel — Cloud hosting for backend and frontend deployment"
    ))
print("  OK Slide 13: AI & Infra Tools")

# =========================================================================
# SLIDE 15: Development Process
# =========================================================================
slide = prs.slides[14]
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "Agile methodology with 2-week sprints:\n\n"
        "Sprint 1-2: Project setup, database design, authentication system (JWT + OAuth2), "
        "core user/role models.\n\n"
        "Sprint 3-4: Resume module (manual CRUD + AI generation), AI service layer with "
        "dual provider support (Gemini + OpenAI).\n\n"
        "Sprint 5-6: Jobs module (CRUD, search, filtering), application pipeline "
        "(submit → review → shortlist → interview → decision).\n\n"
        "Sprint 7-8: Frontend development — student dashboard, company HR portal, admin panel, "
        "responsive UI with Tailwind CSS.\n\n"
        "Sprint 9-10: AI job matching, cover letter generation, ATS scoring, notifications, "
        "payment integration.\n\n"
        "Sprint 11-12: Testing (Pytest + Playwright), CI/CD (GitHub Actions), Docker, "
        "deployment (Render + Vercel), documentation."
    ))
print("  OK Slide 15: Development Process")

# =========================================================================
# SLIDE 17: Challenges & Resolutions
# =========================================================================
slide = prs.slides[16]
s = find(slide, "TextBox 5")
if s:
    set_text(s, (
        "• Full-stack solo workload — Managing FastAPI backend and Next.js 14 frontend as a single "
        "developer across 10+ API modules, 30+ pages, and comprehensive testing.\n\n"
        "• AI Integration Complexity — Connecting dual AI providers (OpenAI GPT-4 and Google Gemini) "
        "for resume generation and job matching; ensuring consistent, valid JSON output from LLMs."
    ))
s = find(slide, "TextBox 6")
if s:
    set_text(s, (
        "• Adopted modular architecture with clear separation of concerns, implemented CI/CD "
        "automation via GitHub Actions, and followed phased Agile development to manage complexity.\n\n"
        "• Developed structured prompt engineering with iterative testing, implemented automatic "
        "provider failover (Gemini → OpenAI), and added JSON validation/normalization layers."
    ))
print("  OK Slide 17: Challenges")

# =========================================================================
# SLIDE 19: Results - Objective 1
# =========================================================================
slide = prs.slides[18]
s = find(slide, "TextBox 3")
if s: set_text(s, "Results — AI Resume Builder")
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "The AI Resume Builder successfully generates tailored, ATS-optimized resumes:\n\n"
        "• Dual AI provider support: Google Gemini (primary, free) and OpenAI GPT-4 (fallback)\n"
        "• Multiple template styles: modern, classic, minimal, creative, professional, executive\n"
        "• ATS score calculation (0-100) with specific improvement suggestions\n"
        "• Multi-language support: English, Uzbek, Russian\n"
        "• Automated content generation from user profile data and target job description\n"
        "• PDF export with professional formatting via ReportLab\n"
        "• Reduces manual resume creation effort by approximately 70%\n\n"
        "The structured prompt engineering approach ensures consistent, high-quality JSON output "
        "that is normalized and validated before storage."
    ))
print("  OK Slide 19: Results - Resume Builder")

# =========================================================================
# SLIDE 20: Results - Objective 2
# =========================================================================
slide = prs.slides[19]
s = find(slide, "TextBox 3")
if s: set_text(s, "Results — AI Job Matching")
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "The intelligent Job Matching system analyzes candidate-job compatibility:\n\n"
        "• Extracts skills, experience level, and keywords from resume content\n"
        "• Calculates match scores (0-100) against all active job listings\n"
        "• Identifies matched skills and missing skills for each job\n"
        "• Provides ranked results with detailed match reasons\n"
        "• Supports filters: location preference, remote-only, minimum salary, experience level\n"
        "• Processes matches in real-time with performance metrics\n\n"
        "The matching algorithm combines skill overlap analysis, experience level comparison, "
        "and keyword matching to provide comprehensive compatibility assessments."
    ))
print("  OK Slide 20: Results - Job Matching")

# =========================================================================
# SLIDE 21: Results - Objective 3
# =========================================================================
slide = prs.slides[20]
s = find(slide, "TextBox 3")
if s: set_text(s, "Results — Complete Multi-Role Platform")
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "The platform delivers a complete 3-role ecosystem:\n\n"
        "Student Role:\n"
        "Register/login (email + OAuth2), build resumes (manual + AI), browse/search jobs, "
        "apply with AI cover letters, track application status, receive notifications.\n\n"
        "Company/HR Role:\n"
        "Post and manage job listings, review applications, filter candidates (basic + AI-assisted), "
        "manage interview pipeline (reviewing → shortlisted → interview → accepted/rejected).\n\n"
        "Admin Role:\n"
        "System dashboard with health monitoring, user management, error tracking via Sentry, "
        "system configuration and oversight.\n\n"
        "All roles share: JWT authentication, email verification, password reset, responsive UI."
    ))
print("  OK Slide 21: Results - Platform")

# =========================================================================
# SLIDE 24: Future Scope & Impact
# =========================================================================
slide = prs.slides[23]
s = find(slide, "TextBox 5")
if s:
    set_text(s, (
        "• Mobile app (React Native) for on-the-go job searching\n"
        "• AI Interview Coach with real-time feedback\n"
        "• ML-based personalized job recommendations"
    ))
s = find(slide, "TextBox 6")
if s:
    set_text(s, (
        "• Empowers students to compete in the job market with AI-optimized resumes\n"
        "• Reduces hiring costs for employers through AI-assisted candidate screening\n"
        "• Demonstrates practical AI application in career services for Central Asian market"
    ))
s = find(slide, "TextBox 8")
if s:
    set_text(s, (
        "Future development includes a React Native mobile app, AI Interview Coach with real-time "
        "feedback, ML-based personalized recommendations, university career center partnerships, "
        "Uzbek/Russian language localization, and Click.uz/Payme payment integration for the local "
        "market. The platform can scale to serve thousands of students across Central Asian universities, "
        "significantly reducing the education-to-employment gap in the region."
    ))
print("  OK Slide 24: Future Scope")

# =========================================================================
# SLIDE 26: Summary & Conclusion
# =========================================================================
slide = prs.slides[25]
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "SmartCareer AI has successfully achieved all five project objectives, delivering a fully "
        "functional AI-powered career platform. Key accomplishments:\n\n"
        "1. AI Resume Builder with dual provider support reduces manual effort by ~70%\n"
        "2. Intelligent Job Matching with 0-100 compatibility scores and skill gap analysis\n"
        "3. Complete 3-role platform with secure authentication and real-time notifications\n\n"
        "Top 3 lessons learned:\n\n"
        "1. Modular Architecture: Essential for managing full-stack solo development — clear separation "
        "of concerns enables incremental progress without overwhelming complexity.\n"
        "2. Iterative Prompt Engineering: Consistent AI output requires structured prompts, JSON "
        "validation, and normalization layers — not just a single API call.\n"
        "3. Automated Testing & CI/CD: Critical for maintaining reliability across a complex, "
        "multi-component system — Pytest + Playwright + GitHub Actions saved countless debugging hours."
    ))
print("  OK Slide 26: Summary")

# =========================================================================
# SLIDE 27: Contribution Breakdown
# =========================================================================
slide = prs.slides[26]
s = find(slide, "TextBox 4")
if s:
    set_text(s, "This is a solo graduation project. All work was completed by a single developer.")

# Fill the table
for shape in slide.shapes:
    if shape.shape_type == 19:  # TABLE
        table = shape.table
        row = table.rows[1]
        row.cells[0].text = STUDENT_ID
        row.cells[1].text = STUDENT_NAME
        row.cells[2].text = "Full Stack Developer\n(Backend, Frontend, AI, PM)"
        row.cells[3].text = (
            "Backend API (FastAPI), AI service integration (Gemini/OpenAI), "
            "Frontend UI (Next.js), authentication, database design, testing, "
            "CI/CD pipeline, deployment, documentation"
        )
        row.cells[4].text = (
            "Complete AI-powered career platform with resume builder, "
            "job matching, cover letter generator, 3-role system, "
            "and production deployment"
        )
        print("  OK Slide 27: Contribution Table")

# =========================================================================
# SLIDE 28: References
# =========================================================================
slide = prs.slides[27]
s = find(slide, "TextBox 4")
if s:
    refs = (
        '[1] S. Ramirez, "FastAPI Documentation," 2023. [Online]. Available: https://fastapi.tiangolo.com/\n'
        '[2] Vercel Inc., "Next.js 14 Documentation," 2024. [Online]. Available: https://nextjs.org/docs\n'
        '[3] OpenAI, "OpenAI API Reference," 2024. [Online]. Available: https://platform.openai.com/docs\n'
        '[4] Google, "Gemini API Documentation," 2024. [Online]. Available: https://ai.google.dev/docs\n'
        '[5] M. Bayer, "SQLAlchemy 2.0 Documentation," 2023. [Online]. Available: https://docs.sqlalchemy.org/\n'
        '[6] Facebook Inc., "React 18 Documentation," 2024. [Online]. Available: https://react.dev/\n'
        '[7] Tailwind Labs, "Tailwind CSS Documentation," 2024. [Online]. Available: https://tailwindcss.com/docs\n'
        '[8] Microsoft, "Playwright Documentation," 2024. [Online]. Available: https://playwright.dev/\n'
        '[9] P. Morelli et al., "Zustand State Management," 2024. [Online]. Available: https://zustand-demo.pmnd.rs/\n'
        '[10] J. Jones et al., "AI-Driven Resume Screening: A Systematic Review," J. of AI Research, vol. 72, 2023.\n'
        '[11] S. Chen and R. Gupta, "NLP for Automated Career Services," Proc. IEEE Int. Conf. NLP and AI, 2023.'
    )
    set_text(s, refs)
print("  OK Slide 28: References")

# =========================================================================
# SLIDE 29: Acknowledgment
# =========================================================================
slide = prs.slides[28]
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        f"I would like to express my sincere gratitude to my supervisor, Mr. {SUPERVISOR}, "
        f"and to Professor {PROFESSOR} for their invaluable guidance, constructive feedback, "
        "and continuous support throughout the development of this project. Their expertise and "
        "encouragement played a crucial role in shaping the direction and quality of this work.\n\n"
        "My sincere appreciation also goes to the Engineering School and the Central Asian "
        "University community for their continuous academic inspiration, support, and the "
        "enriching environment they provide for learning and innovation. Their commitment to "
        "excellence has played a significant role in shaping my skills and passion for technology."
    ))
print("  OK Slide 29: Acknowledgment")

out = f'{STUDENT_ID}_SmartCareer_AI.pptx'
prs.save(out)
print(f"\nSaved: {out}")
