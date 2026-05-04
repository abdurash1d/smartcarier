#!/usr/bin/env python3
"""Quality review fixes for both files."""
from pptx import Presentation
from pptx.oxml.ns import qn
import copy

STUDENT_ID = "210004"

def set_text(shape, text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        tf.paragraphs[i]._p.getparent().remove(tf.paragraphs[i]._p)
    para = tf.paragraphs[0]
    for run in list(para.runs):
        run._r.getparent().remove(run._r)
    lines = text.split('\n')
    for li, line in enumerate(lines):
        if li == 0:
            r = para.add_run()
            r.text = line
        else:
            new_p = copy.deepcopy(para._p)
            for old_r in new_p.findall(qn('a:r')):
                new_p.remove(old_r)
            new_r = copy.deepcopy(para._p.findall(qn('a:r'))[0])
            new_r.find(qn('a:t')).text = line
            new_p.append(new_r)
            tf._txBody.append(new_p)

def find(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


# =========================================================================
# POSTER REFINEMENTS
# =========================================================================
print("REFINING POSTER...")
prs = Presentation(f'{STUDENT_ID}_Poster.pptx')
slide = prs.slides[0]

# Refined INTRODUCTION (more concise, sharper opening)
s = find(slide, "TextBox 31")
if s:
    set_text(s, (
        "INTRODUCTION\n\n"
        "The transition from education to employment is one of the most critical "
        "challenges faced by university students and recent graduates. Despite the "
        "abundance of online job portals, many young job seekers struggle to create "
        "professional, ATS-optimized resumes and to match their skills to suitable "
        "opportunities. At the same time, employers face difficulties in efficiently "
        "filtering large applicant pools to identify the most qualified candidates. "
        "This bidirectional gap highlights the need for an intelligent, AI-driven "
        "platform capable of streamlining the entire career development pipeline.\n\n"
        "SmartCareer AI addresses this challenge by integrating Google Gemini and "
        "OpenAI GPT-4 to automate resume generation, analyze resume-to-job compatibility, "
        "generate tailored cover letters, and facilitate intelligent job matching — "
        "empowering students to compete more effectively in today's job market while "
        "providing employers with AI-assisted candidate screening tools."
    ))
    print("  Refined: Introduction (clearer two-paragraph structure)")

# Refined METHODOLOGY (cleaner bullet formatting)
s = find(slide, "TextBox 35")
if s:
    set_text(s, (
        "METHODOLOGY\n\n"
        "The project followed an Agile development methodology with iterative two-week "
        "sprints. The system architecture is organized into four cohesive layers:\n\n"
        "• Backend Layer — FastAPI (Python) with versioned RESTful endpoints, "
        "SQLAlchemy ORM over PostgreSQL, Alembic migrations, and JWT + OAuth2 (Google) "
        "authentication for secure multi-role access.\n\n"
        "• AI Service Layer — A dual-provider abstraction supports Google Gemini "
        "(primary, free) and OpenAI GPT-4 (fallback) with automatic failover. Structured "
        "prompts and JSON-validated responses ensure consistent output quality.\n\n"
        "• Frontend Layer — Next.js 14 with React 18 and TypeScript delivers server-side "
        "rendering for performance and SEO. Zustand manages state; Tailwind CSS and "
        "Radix UI provide a modern, accessible interface.\n\n"
        "• Quality & Deployment — Pytest for backend unit tests, Playwright for "
        "end-to-end testing, GitHub Actions for CI/CD, and Docker-based deployment "
        "to Render and Vercel."
    ))
    print("  Refined: Methodology (clear bullet hierarchy)")

# Refined CONCLUSION (clearer lessons-learned structure)
s = find(slide, "TextBox 50")
if s:
    set_text(s, (
        "CONCLUSION\n\n"
        "This project has successfully achieved all five of its research objectives, "
        "delivering a fully functional AI-powered career platform that bridges the gap "
        "between job seekers and employers. SmartCareer AI has evolved beyond a simple "
        "job portal into an intelligent career development tool that leverages cutting-edge "
        "language models to reduce friction in resume creation and job matching. The "
        "dual-provider AI architecture (Gemini + OpenAI) ensures both reliability and "
        "cost-effectiveness, while the modular codebase and comprehensive testing "
        "infrastructure provide a solid foundation for future development.\n\n"
        "Three key lessons emerged from the work: (1) modular architecture is essential "
        "for managing full-stack solo development; (2) iterative prompt engineering — "
        "with JSON validation and normalization layers — is critical for consistent AI "
        "output; and (3) automated testing combined with CI/CD pipelines is indispensable "
        "for maintaining reliability across complex, multi-component systems."
    ))
    print("  Refined: Conclusion (numbered lessons learned)")

# Refined RESULTS (slight polish)
s = find(slide, "TextBox 36")
if s:
    set_text(s, (
        "RESULTS\n\n"
        "The AI Resume Builder reliably generates tailored, ATS-optimized resumes using "
        "dual AI providers (Gemini and GPT-4), reducing manual resume creation effort "
        "by approximately 70%. The intelligent Job Matching system analyzes candidate "
        "profiles against job requirements and produces compatibility scores (0–100), "
        "matched-skill lists, missing-skill gaps, and human-readable match reasons. "
        "The platform delivers a complete three-role ecosystem: students build resumes, "
        "search and apply for jobs, and track application status; companies post jobs, "
        "review applicants, and filter candidates with AI assistance; administrators "
        "monitor system health, users, and errors through a dedicated dashboard. The "
        "Cover Letter Generator produces personalized, job-specific letters aligned to "
        "the candidate's resume. The architecture proved robust under testing, with "
        "Sentry-based error monitoring, Redis-backed rate limiting, and GZip compression "
        "ensuring production-grade performance and reliability."
    ))
    print("  Refined: Results (more measurable and structured)")

# Refined FUTURE WORK (sharper academic tone)
s = find(slide, "TextBox 52")
if s:
    set_text(s, (
        "RECOMMENDATION AND FUTURE WORK\n\n"
        "While this project delivered a complete web-based platform, several promising "
        "directions remain for future research and development:\n\n"
        "1. Mobile Application — A React Native client to enable on-the-go job search "
        "and application tracking.\n"
        "2. AI Interview Coach — A conversational AI module that conducts mock interviews "
        "and provides real-time feedback on responses.\n"
        "3. Machine Learning Recommendations — Behavioral models trained on user "
        "interaction data for personalized job recommendations beyond keyword matching.\n"
        "4. University Integration — Partnerships with career centers to incorporate "
        "grant search and motivation-letter generation.\n"
        "5. Localization — Full Uzbek and Russian language support for the Central "
        "Asian job market.\n"
        "6. Local Payments — Integration of Click.uz and Payme alongside the existing "
        "Stripe gateway."
    ))
    print("  Refined: Future Work (academic tone)")

# Refined PROJECT OUTPUT (drop quotation marks remnants)
s = find(slide, "TextBox 57")
if s:
    set_text(s, (
        "PROJECT OUTPUT\n\n"
        "The screenshots below illustrate the SmartCareer AI web platform interface, "
        "including the AI Resume Builder, Job Matching dashboard, Application Tracker, "
        "Company HR portal, and Admin dashboard — demonstrating the full feature set "
        "across all three user roles."
    ))
    print("  Refined: Project Output")

# Refined ACKNOWLEDGMENT (cleaner academic register)
s = find(slide, "TextBox 49")
if s:
    set_text(s, (
        "ACKNOWLEDGMENT\n\n"
        "I would like to express my sincere gratitude to my supervisor, Mr. Farhod "
        "Mahmudxo'jayev, and to Professor Eugene Castro, for their invaluable guidance, "
        "constructive feedback, and continuous support throughout the development of "
        "this project. My sincere appreciation also extends to the Engineering School "
        "and the wider Central Asian University community for fostering an enriching "
        "academic environment that has played a central role in shaping my skills and "
        "passion for technology."
    ))
    print("  Refined: Acknowledgment")

prs.save(f'{STUDENT_ID}_Poster.pptx')
print(f"  Saved: {STUDENT_ID}_Poster.pptx\n")


# =========================================================================
# PRESENTATION REFINEMENTS
# =========================================================================
print("REFINING PRESENTATION...")
prs = Presentation(f'{STUDENT_ID}_SmartCareer_AI.pptx')

# Slide 8: Fix redundant "/applications — Applications,"
slide = prs.slides[7]
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "The backend exposes ten versioned API modules under /api/v1/:\n\n"
        "• /auth — Registration, login, email verification, password reset, OAuth2 (Google)\n"
        "• /users — User profile management and role-based access\n"
        "• /resumes — Manual resume CRUD and AI resume generation with ATS scoring\n"
        "• /jobs — Job posting CRUD, search, filtering, publish/unpublish workflows\n"
        "• /applications — Job applications and status tracking pipeline "
        "(reviewing → shortlisted → interview → accepted/rejected)\n"
        "• /ai — AI Resume Generator, Resume Analyzer, Cover Letter Generator, Job Matcher\n"
        "• /admin — System health, user management, and error monitoring dashboard\n"
        "• /payments — Stripe subscription and billing management\n"
        "• /notifications — Real-time in-app notification system\n"
        "• /saved-searches — Persistent search filters for users\n\n"
        "All endpoints provide auto-generated Swagger/ReDoc documentation, Pydantic "
        "request validation, structured error envelopes with request IDs, and security "
        "headers (HSTS, CSP, X-Frame-Options, XSS protection)."
    ))
    print("  Refined: Slide 8 — API Structure")

# Slide 17: Fix title trailing space
slide = prs.slides[16]
s = find(slide, "TextBox 3")
if s and s.text.strip() == "Challenges":
    set_text(s, "Challenges")
    print("  Refined: Slide 17 — Removed trailing space in 'Challenges'")

# Slide 1: Cleaner subtitle / student info layout
slide = prs.slides[0]
s = find(slide, "TextBox 7")
if s:
    set_text(s, (
        "Graduation Project\n\n"
        "Djumabaev Abdurashid  |  Student ID: 210004\n"
        "BSc in Computer Science\n\n"
        "Supervisor: Mr. Farhod Mahmudxo'jayev\n"
        "Professor: Eugene Castro"
    ))
    print("  Refined: Slide 1 — Cleaner student info layout")

# Slide 26: Restructure summary for better readability
slide = prs.slides[25]
s = find(slide, "TextBox 4")
if s:
    set_text(s, (
        "SmartCareer AI has successfully achieved all five project objectives, delivering "
        "a fully functional AI-powered career platform. Key accomplishments include:\n\n"
        "• AI Resume Builder with dual provider support, reducing manual effort by ~70%\n"
        "• Intelligent Job Matching with 0–100 compatibility scores and skill gap analysis\n"
        "• Complete three-role platform with secure authentication and real-time notifications\n\n"
        "Top three lessons learned:\n\n"
        "1. Modular Architecture — Essential for managing full-stack solo development; "
        "clear separation of concerns enables incremental progress.\n\n"
        "2. Iterative Prompt Engineering — Consistent AI output requires structured "
        "prompts, JSON validation, and normalization layers, not just a single API call.\n\n"
        "3. Automated Testing & CI/CD — Critical for reliability across complex systems; "
        "Pytest, Playwright, and GitHub Actions saved countless debugging hours."
    ))
    print("  Refined: Slide 26 — Summary structure")

# Slide 5: Realistic timeline (clarify scope)
slide = prs.slides[4]
s = find(slide, "TextBox 5")
if s:
    set_text(s, (
        "The project was executed over approximately 14 weeks using Agile methodology "
        "with two-week sprints:\n\n"
        "Phase 1 — Discovery & Design (Weeks 1–2): Requirements analysis, system "
        "architecture, database schema design, and project setup.\n\n"
        "Phase 2 — Backend Foundation (Weeks 3–5): FastAPI setup, JWT + OAuth2 "
        "authentication, SQLAlchemy models, and core API endpoints.\n\n"
        "Phase 3 — AI Integration (Weeks 6–8): OpenAI and Gemini service implementation, "
        "resume generation, job matching algorithm, ATS scoring, and cover letter generator.\n\n"
        "Phase 4 — Frontend Development (Weeks 9–11): Next.js 14 setup, student dashboard, "
        "company HR portal, admin panel, and responsive UI with Tailwind CSS.\n\n"
        "Phase 5 — Testing & Deployment (Weeks 12–13): Pytest unit tests, Playwright E2E "
        "tests, GitHub Actions CI/CD, Docker containerization, and deployment to Render and Vercel.\n\n"
        "Phase 6 — Finalization (Week 14): Final testing, documentation, and presentation preparation."
    ))
    print("  Refined: Slide 5 — Timeline with named phases")

# Slide 3: Cleaner motivation
slide = prs.slides[2]
s = find(slide, "TextBox 5")
if s:
    set_text(s, (
        "Industry research indicates that over 75% of resumes are rejected by Applicant "
        "Tracking Systems (ATS) before reaching a human recruiter, while junior candidates "
        "often spend 40+ hours per month on job applications. SmartCareer AI directly "
        "addresses this inefficiency by leveraging artificial intelligence to automate "
        "resume creation, deliver intelligent job matching, and streamline the entire "
        "application workflow for both candidates and employers."
    ))
    print("  Refined: Slide 3 — Motivation")

prs.save(f'{STUDENT_ID}_SmartCareer_AI.pptx')
print(f"  Saved: {STUDENT_ID}_SmartCareer_AI.pptx")

print("\n" + "=" * 60)
print("QUALITY REVIEW COMPLETE")
print("=" * 60)
